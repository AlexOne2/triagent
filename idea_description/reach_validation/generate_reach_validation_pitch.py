from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "idea_description" / "reach_validation" / "triagent_reach_validation_pitch.pptx"
SCREENSHOT = ROOT / "Dashboard.png"


COLORS = {
    "bg": RGBColor(248, 250, 252),
    "surface": RGBColor(255, 255, 255),
    "surface_alt": RGBColor(239, 246, 251),
    "surface_warm": RGBColor(255, 247, 237),
    "surface_green": RGBColor(236, 253, 245),
    "surface_amber": RGBColor(255, 247, 214),
    "border": RGBColor(215, 227, 234),
    "text": RGBColor(15, 23, 42),
    "muted": RGBColor(71, 85, 105),
    "brand": RGBColor(29, 95, 122),
    "brand_dark": RGBColor(23, 40, 71),
    "amber": RGBColor(181, 107, 0),
    "green": RGBColor(22, 101, 52),
    "red": RGBColor(220, 38, 38),
}


def set_background(slide, color_key: str = "bg") -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[color_key]


def add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: str = "surface",
    line: str = "border",
    radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(radius_shape, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    shape.line.width = Pt(1)
    return shape


def add_divider(
    slide,
    left: float,
    top: float,
    width: float,
    *,
    height: float = 0.03,
    color: str = "border",
) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[color]
    bar.line.fill.background()


def add_vdivider(
    slide,
    left: float,
    top: float,
    height: float,
    *,
    width: float = 0.03,
    color: str = "border",
) -> None:
    add_divider(slide, left, top, width, height=height, color=color)


def style_text_frame(text_frame, *, margin=0.12) -> None:
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    color: str = "text",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    style_text_frame(box.text_frame)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_label(slide, text: str, *, left: float, top: float, width: float = 2.4, color: str = "brand") -> None:
    add_textbox(slide, left, top, width, 0.25, text.upper(), font_size=10, color=color, bold=True)


def add_body_lines(
    shape,
    lines: list[str],
    *,
    font_size: int = 15,
    color: str = "muted",
    bullet: bool = False,
    spacing_after: int = 6,
) -> None:
    tf = shape.text_frame
    tf.clear()
    style_text_frame(tf)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line if not bullet else f"• {line}"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing_after)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.color.rgb = COLORS[color]


def add_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    lines: list[str],
    fill: str = "surface",
    line: str = "border",
    title_color: str = "text",
    body_color: str = "muted",
) -> None:
    accent_color = line if line != "border" else "brand"
    add_divider(slide, left, top, width, color=accent_color)
    card = slide.shapes.add_textbox(Inches(left), Inches(top + 0.1), Inches(width), Inches(height - 0.1))
    tf = card.text_frame
    tf.clear()
    style_text_frame(tf, margin=0.18)

    p = tf.paragraphs[0]
    p.text = title
    p.space_after = Pt(8)
    run = p.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(19)
    run.font.bold = True
    run.font.color.rgb = COLORS[title_color]

    for line_text in lines:
        p = tf.add_paragraph()
        p.text = f"• {line_text}"
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(13)
        run.font.color.rgb = COLORS[body_color]


def add_stat(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    number: str,
    title: str,
    subtitle: str,
) -> None:
    add_divider(slide, left, top, width, color="brand")
    box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.08), Inches(width), Inches(1.25))
    tf = box.text_frame
    tf.clear()
    style_text_frame(tf, margin=0.05)
    p = tf.paragraphs[0]
    p.text = number
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = COLORS["brand_dark"]
    p = tf.add_paragraph()
    p.text = title
    p.space_after = Pt(4)
    run = p.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS["text"]
    p = tf.add_paragraph()
    p.text = subtitle
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS["muted"]


def add_footer(slide, index: int) -> None:
    add_textbox(
        slide,
        10.8,
        6.95,
        2.0,
        0.25,
        f"Triagent · {index}/6",
        font_size=9,
        color="muted",
        align=PP_ALIGN.RIGHT,
    )


def slide_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_label(slide, "Problem", left=0.6, top=0.35, width=1.4)
    add_textbox(
        slide,
        0.6,
        0.75,
        9.8,
        0.9,
        "Phishing triage is still too manual.",
        font_size=28,
        bold=True,
        color="text",
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        0.6,
        1.6,
        11.4,
        0.8,
        "User-reported suspicious emails still turn into scattered checks across mail clients, URL tools, sandboxes, SIEM searches, and tickets.",
        font_size=17,
        color="muted",
    )

    add_card(
        slide,
        left=0.6,
        top=2.65,
        width=3.95,
        height=2.1,
        title="Repetitive analysis",
        lines=[
            "Analysts repeat the same evidence-gathering steps for each reported email.",
            "Similar emails are often handled one by one instead of as one wave.",
        ],
        fill="surface",
    )
    add_card(
        slide,
        left=4.7,
        top=2.65,
        width=3.95,
        height=2.1,
        title="Defensible decisions",
        lines=[
            "The job is not only a verdict; it is rationale, artifacts, and clean closure.",
            "Documentation and auditability add real workload.",
        ],
        fill="surface",
    )
    add_card(
        slide,
        left=8.8,
        top=2.65,
        width=3.95,
        height=2.1,
        title="Regulated friction",
        lines=[
            "Cloud-first tooling is harder to adopt when email content should stay local.",
            "A focused on-prem workflow is easier to discuss than a broad suite replacement.",
        ],
        fill="surface",
    )

    add_divider(slide, 0.6, 5.18, 12.15, color="brand")
    add_textbox(
        slide,
        0.6,
        5.28,
        12.0,
        0.55,
        "Hypothesis: the sharper problem is not better phishing detection alone. It is faster evidence gathering, categorization, and documentation for reported email.",
        font_size=15,
        color="brand_dark",
        bold=True,
    )

    add_textbox(
        slide,
        0.6,
        6.25,
        6.0,
        0.25,
        "Founder context: Alexander Huelsmann · MSc thesis on AI-assisted phishing triage",
        font_size=10,
        color="muted",
    )
    add_footer(slide, 1)


def slide_validation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_label(slide, "Validation", left=0.6, top=0.35, width=1.8)
    add_textbox(
        slide,
        0.6,
        0.75,
        7.6,
        0.75,
        "What I validated so far.",
        font_size=28,
        bold=True,
        color="text",
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        0.6,
        1.45,
        11.5,
        0.55,
        "Current evidence comes from interviews, follow-up chats, and AI-assisted Reddit workflow research. Buyer validation is the next open risk.",
        font_size=16,
        color="muted",
    )

    metrics = [
        ("7", "Interviews", "Direct analyst and operator workflow conversations completed"),
        ("5", "Follow-up chats", "Additional validation conversations across the same problem area"),
        ("AI", "Reddit research", "AI-assisted scraping and synthesis of analyst pain and workflow patterns"),
    ]
    for idx, (number, title, subtitle) in enumerate(metrics):
        add_stat(slide, left=0.6 + idx * 4.1, top=2.25, width=3.8, number=number, title=title, subtitle=subtitle)
        if idx < len(metrics) - 1:
            add_vdivider(slide, 4.5 + idx * 4.1, 2.25, 1.45)

    add_divider(slide, 0.6, 4.1, 7.55, color="brand")
    learnings = slide.shapes.add_textbox(Inches(0.6), Inches(4.18), Inches(7.55), Inches(2.1))
    tf = learnings.text_frame
    tf.clear()
    style_text_frame(tf, margin=0.02)
    p = tf.paragraphs[0]
    p.text = "Recurring learnings"
    run = p.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS["text"]
    for line in [
        "Reported phishing handling is a daily analyst task and still largely manual.",
        "On-prem / privacy constraints matter; a smaller focused tool is easier to discuss than a broad cloud suite.",
        "URL shorteners and attachment edge cases still require analyst review. Black-box automation is not enough.",
        "The wedge shifted from better detection to better analyst workflow.",
    ]:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.space_after = Pt(5)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(13)
        run.font.color.rgb = COLORS["muted"]

    add_divider(slide, 8.45, 4.1, 4.3, color="amber")
    tf_box = slide.shapes.add_textbox(Inches(8.45), Inches(4.18), Inches(4.3), Inches(2.1))
    tf = tf_box.text_frame
    tf.clear()
    style_text_frame(tf, margin=0.02)
    p = tf.paragraphs[0]
    p.text = "Open validation question"
    run = p.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS["amber"]
    for line in [
        "The analyst pain is clear, but buyer sponsorship still needs proof.",
        "Next step: validate whether SOC managers or MSSPs will back a pilot.",
        "Target outcome: 1-3 design partners with real workflow access.",
    ]:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(13)
        run.font.color.rgb = COLORS["brand_dark"]

    add_footer(slide, 2)


def slide_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_label(slide, "Solution", left=0.6, top=0.35, width=1.6)
    add_textbox(
        slide,
        0.6,
        0.75,
        8.6,
        0.75,
        "Triagent: analyst-in-the-loop phishing triage.",
        font_size=28,
        bold=True,
        color="text",
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        0.6,
        1.45,
        11.6,
        0.55,
        "Triagent does not replace the analyst. It removes repetitive evidence collection and packages the decision.",
        font_size=16,
        color="muted",
    )

    steps = [
        ("01", "Ingest reported email", "Start from the original message or the Outlook add-in."),
        ("02", "Auto-categorize", "Separate likely benign reports from analyst-worthy cases."),
        ("03", "Enrich evidence", "Headers, URLs, attachments, auth signals, and context are collected."),
        ("04", "Assist draft + analyst review", "Triagent drafts a resolution, but the analyst approves the decision."),
        ("05", "Export the case", "Generate evidence reports, IOCs, and audit-ready history."),
    ]
    left = 0.6
    top = 2.35
    width = 2.35
    height = 2.65
    gap = 0.12
    for idx, (num, title, body) in enumerate(steps):
        x = left + idx * (width + gap)
        add_divider(slide, x, top, width, color="brand")
        card = slide.shapes.add_textbox(Inches(x), Inches(top + 0.08), Inches(width), Inches(height - 0.08))
        tf = card.text_frame
        tf.clear()
        style_text_frame(tf, margin=0.04)
        p = tf.paragraphs[0]
        p.text = num
        run = p.runs[0]
        run.font.name = "Aptos Display"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = COLORS["brand"]
        p = tf.add_paragraph()
        p.text = title
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Aptos Display"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS["text"]
        p = tf.add_paragraph()
        p.text = body
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS["muted"]

    add_divider(slide, 0.6, 5.45, 12.15, color="brand")
    add_textbox(
        slide,
        0.6,
        5.54,
        12.1,
        0.42,
        "Automatic categorization protects analyst time; evidence-first review makes the remaining cases faster and more defensible.",
        font_size=15,
        color="brand_dark",
        bold=True,
    )
    add_footer(slide, 3)


def slide_progress(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_label(slide, "Progress", left=0.6, top=0.35, width=1.7)
    add_textbox(
        slide,
        0.6,
        0.75,
        7.8,
        0.75,
        "From idea to working prototype.",
        font_size=28,
        bold=True,
        color="text",
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        0.6,
        1.45,
        11.6,
        0.55,
        "The goal of the last weeks was not polish; it was to turn the concept into something demoable and testable.",
        font_size=16,
        color="muted",
    )

    if SCREENSHOT.exists():
        slide.shapes.add_picture(str(SCREENSHOT), Inches(0.6), Inches(2.05), width=Inches(7.0))

    add_vdivider(slide, 7.7, 2.05, 3.55)
    card = slide.shapes.add_textbox(Inches(8.0), Inches(2.05), Inches(4.75), Inches(3.55))
    tf = card.text_frame
    tf.clear()
    style_text_frame(tf, margin=0.02)
    p = tf.paragraphs[0]
    p.text = "Already working today"
    run = p.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS["text"]
    for line in [
        ".eml / .msg intake and Outlook add-in path",
        "In-tray, report workspace, and analyst resolution",
        "URL, attachment, auth, transmission, and raw-source analysis",
        "Local assist draft with flagged artifacts",
        "PDF / IOC export and tamper-evident audit log",
        "Synthetic corpus plus one-command demo reset",
    ]:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.space_after = Pt(5)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS["muted"]

    add_divider(slide, 0.6, 6.05, 12.15, color="amber")
    add_textbox(
        slide,
        0.6,
        6.14,
        12.0,
        0.4,
        "Current status: validation-grade prototype, not yet an enterprise-ready platform. Next goal: test one complete workflow on real or semi-real cases.",
        font_size=14,
        color="brand_dark",
        bold=True,
    )
    add_footer(slide, 4)


def slide_business(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_label(slide, "Market / Model / USP", left=0.6, top=0.35, width=3.2)
    add_textbox(
        slide,
        0.6,
        0.75,
        8.5,
        0.75,
        "Beachhead first, not broad cyber.",
        font_size=28,
        bold=True,
        color="text",
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        0.6,
        1.45,
        11.4,
        0.55,
        "The initial market is narrow on purpose: phishing triage workflows in regulated SOCs and MSSPs.",
        font_size=16,
        color="muted",
    )

    add_card(
        slide,
        left=0.6,
        top=2.2,
        width=3.9,
        height=3.45,
        title="Beachhead customer",
        lines=[
            "Internal SOC / SecOps teams with recurring phishing mailbox volume",
            "MSSPs serving regulated or privacy-sensitive clients",
            "Likely champion: SOC manager, MSSP service lead, or CISO",
        ],
    )
    add_card(
        slide,
        left=4.72,
        top=2.2,
        width=3.9,
        height=3.45,
        title="Business model",
        lines=[
            "Design-partner pilots first",
            "Then annual on-prem licence or managed deployment fee",
            "Focus on workflow pull before optimizing pricing",
        ],
        fill="surface_alt",
        line="surface_alt",
    )
    add_card(
        slide,
        left=8.84,
        top=2.2,
        width=3.9,
        height=3.45,
        title="Alternatives / USP",
        lines=[
            "Alternatives today: manual workflow, broad suites, SOAR, phishing-analysis tools",
            "USP: on-prem, workflow-first, evidence-report-first",
            "Campaign clustering is part of the product direction",
        ],
    )

    add_divider(slide, 0.6, 6.0, 12.15, color="green")
    add_textbox(
        slide,
        0.6,
        6.09,
        12.0,
        0.38,
        "Biggest commercial question: will buyers pay for this narrower workflow wedge, or treat it as just a feature?",
        font_size=14,
        color="green",
        bold=True,
    )
    add_footer(slide, 5)


def slide_ask(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_label(slide, "Roadmap / CTA", left=0.6, top=0.35, width=2.2)
    add_textbox(
        slide,
        0.6,
        0.75,
        8.0,
        0.75,
        "Next milestone: 3 design partners.",
        font_size=28,
        bold=True,
        color="text",
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        0.6,
        1.45,
        11.2,
        0.55,
        "This is a progress pitch: I am not claiming PMF yet. I am asking for access to test the workflow in real settings.",
        font_size=16,
        color="muted",
    )

    steps = [
        ("Next 2 weeks", "Polish one complete reported-email triage workflow for repeatable demos."),
        ("Next 4 weeks", "Run shadow pilots on real or semi-real reported phishing cases."),
        ("Next 8 weeks", "Measure time saved, report quality, and willingness to pay."),
    ]
    for idx, (label, body) in enumerate(steps):
        y = 2.2 + idx * 1.25
        add_divider(slide, 0.6, y, 6.0, color="brand")
        card = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.08), Inches(6.0), Inches(0.9))
        tf = card.text_frame
        tf.clear()
        style_text_frame(tf, margin=0.02)
        p = tf.paragraphs[0]
        p.text = label
        p.space_after = Pt(4)
        run = p.runs[0]
        run.font.name = "Aptos Display"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS["brand"]
        p = tf.add_paragraph()
        p.text = body
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS["muted"]

    add_vdivider(slide, 6.75, 2.2, 3.95)
    ask = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.75), Inches(3.95))
    tf = ask.text_frame
    tf.clear()
    style_text_frame(tf, margin=0.02)
    p = tf.paragraphs[0]
    p.text = "What I’m looking for"
    run = p.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS["brand_dark"]
    for line in [
        "Introductions to SOC managers, MSSP operators, and IT security leads",
        "30-minute workflow interviews",
        "1-3 design-partner or shadow-pilot opportunities",
        "Network and validation support before capital",
    ]:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(13)
        run.font.color.rgb = COLORS["brand_dark"]

    add_divider(slide, 0.6, 6.35, 12.15, color="brand_dark")
    add_textbox(
        slide,
        0.6,
        6.44,
        12.0,
        0.36,
        "Current status: founder-market fit + early analyst validation + working prototype. Next step: buyer access.",
        font_size=14,
        color="brand_dark",
        bold=True,
    )
    add_footer(slide, 6)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    core = prs.core_properties
    core.author = "Alexander Huelsmann / Codex"
    core.title = "Triagent REACH Validation Pitch"
    core.subject = "REACH Startup Validation Lab final pitch"
    core.keywords = "Triagent, phishing triage, validation pitch, REACH"
    core.comments = "Generated from repo-grounded product context and validation notes."
    core.created = datetime.now(timezone.utc)

    slide_problem(prs)
    slide_validation(prs)
    slide_solution(prs)
    slide_progress(prs)
    slide_business(prs)
    slide_ask(prs)
    return prs


def main() -> None:
    prs = build_deck()
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
